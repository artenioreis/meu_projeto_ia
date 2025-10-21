import numpy as np
import nltk
import string
import random
import os
import glob
import re
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# Importações com tratamento de erro robusto
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("⚠️  python-docx não disponível")

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("⚠️  PyPDF2 não disponível")

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("⚠️  openpyxl não disponível")

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx não disponível")

try:
    from odf import text, teletype
    ODF_AVAILABLE = True
except ImportError:
    ODF_AVAILABLE = False
    print("⚠️  odfpy não disponível")

try:
    import speech_recognition as sr
    AUDIO_AVAILABLE = True
except ImportError:
    AUDIO_AVAILABLE = False
    print("⚠️  speechrecognition não disponível")

try:
    from moviepy.editor import VideoFileClip
    VIDEO_AVAILABLE = True
except ImportError:
    VIDEO_AVAILABLE = False
    print("⚠️  moviepy não disponível")

class AdamIAMultimodal:
    def __init__(self):
        """
        Adam Multimodal - IA que aprende com TODOS os tipos de arquivos
        """
        self.respostas = []
        self.fontes_conhecimento = {}
        self.vectorizer = TfidfVectorizer(max_features=5000, stop_words='english')
        self.caracteristicas = None
        
        self.estatisticas = {
            'total_sentencas': 0,
            'total_arquivos': 0,
            'arquivos_processados': [],
            'tipos_processados': {},
            'ultima_atualizacao': None,
            'erros_processamento': []
        }
        
        self._baixar_recursos_nltk()
        self._verificar_dependencias()
    
    def _baixar_recursos_nltk(self):
        """Baixa recursos do NLTK"""
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            print("📥 Baixando recursos NLTK...")
            nltk.download('punkt')
    
    def _verificar_dependencias(self):
        """Verifica quais dependências estão disponíveis"""
        print("🔍 Adam Multimodal - Capacidades:")
        capacidades = []
        if DOCX_AVAILABLE: capacidades.append("📄 Word")
        if PDF_AVAILABLE: capacidades.append("📑 PDF") 
        if EXCEL_AVAILABLE: capacidades.append("📊 Excel")
        if PPTX_AVAILABLE: capacidades.append("🎤 PowerPoint")
        if ODF_AVAILABLE: capacidades.append("📝 OpenDocument")
        if AUDIO_AVAILABLE: capacidades.append("🎵 Áudio")
        if VIDEO_AVAILABLE: capacidades.append("🎬 Vídeo")
        
        if capacidades:
            print(f"✅ {', '.join(capacidades)}")
        else:
            print("⚠️  Nenhuma capacidade adicional disponível")
    
    def preprocessar_texto(self, texto):
        """Limpa e prepara o texto"""
        if not texto:
            return ""
        texto = str(texto)
        texto = texto.translate(str.maketrans('', '', string.punctuation))
        texto = texto.lower()
        texto = re.sub(r'\s+', ' ', texto)
        return texto.strip()
    
    # ========== PROCESSADORES DE ARQUIVOS ==========
    
    def processar_txt(self, arquivo_path):
        """Processa arquivos .txt"""
        try:
            # Tenta diferentes encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(arquivo_path, 'r', encoding=encoding) as f:
                        conteudo = f.read()
                    if conteudo.strip():
                        sentencas = nltk.sent_tokenize(conteudo)
                        return [s for s in sentencas if len(s.strip()) > 5]
                except UnicodeDecodeError:
                    continue
            return []
        except Exception as e:
            self.estatisticas['erros_processamento'].append(f"TXT {os.path.basename(arquivo_path)}: {e}")
            return []
    
    def processar_docx(self, arquivo_path):
        """Processa arquivos .docx"""
        if not DOCX_AVAILABLE:
            return []
        
        try:
            doc = docx.Document(arquivo_path)
            texto_completo = []
            for para in doc.paragraphs:
                if para.text and para.text.strip():
                    texto_completo.append(para.text.strip())
            
            texto = ' '.join(texto_completo)
            if texto:
                sentencas = nltk.sent_tokenize(texto)
                return [s for s in sentencas if len(s.strip()) > 5]
            return []
        except Exception as e:
            self.estatisticas['erros_processamento'].append(f"DOCX {os.path.basename(arquivo_path)}: {e}")
            return []
    
    def processar_pdf(self, arquivo_path):
        """Processa arquivos .pdf"""
        if not PDF_AVAILABLE:
            return []
        
        try:
            texto_completo = []
            with open(arquivo_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for pagina_num, pagina in enumerate(pdf_reader.pages):
                    try:
                        texto = pagina.extract_text()
                        if texto and texto.strip():
                            texto_completo.append(texto.strip())
                    except:
                        continue
            
            texto = ' '.join(texto_completo)
            if texto:
                sentencas = nltk.sent_tokenize(texto)
                return [s for s in sentencas if len(s.strip()) > 5]
            return []
        except Exception as e:
            self.estatisticas['erros_processamento'].append(f"PDF {os.path.basename(arquivo_path)}: {e}")
            return []
    
    def processar_xlsx(self, arquivo_path):
        """Processa arquivos .xlsx"""
        if not EXCEL_AVAILABLE:
            return []
        
        try:
            texto_completo = []
            workbook = openpyxl.load_workbook(arquivo_path, read_only=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows(values_only=True):
                    linha_texto = ' '.join([str(cell) for cell in row if cell is not None and str(cell).strip()])
                    if linha_texto.strip():
                        texto_completo.append(linha_texto.strip())
            
            texto = ' '.join(texto_completo)
            if texto:
                sentencas = nltk.sent_tokenize(texto)
                return [s for s in sentencas if len(s.strip()) > 5]
            return []
        except Exception as e:
            self.estatisticas['erros_processamento'].append(f"XLSX {os.path.basename(arquivo_path)}: {e}")
            return []
    
    def processar_pptx(self, arquivo_path):
        """Processa arquivos .pptx"""
        if not PPTX_AVAILABLE:
            return []
        
        try:
            texto_completo = []
            presentation = pptx.Presentation(arquivo_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        texto_completo.append(shape.text.strip())
            
            texto = ' '.join(texto_completo)
            if texto:
                sentencas = nltk.sent_tokenize(texto)
                return [s for s in sentencas if len(s.strip()) > 5]
            return []
        except Exception as e:
            self.estatisticas['erros_processamento'].append(f"PPTX {os.path.basename(arquivo_path)}: {e}")
            return []
    
    def processar_odt(self, arquivo_path):
        """Processa arquivos .odt"""
        if not ODF_AVAILABLE:
            return []
        
        try:
            from odf.opendocument import load
            doc = load(arquivo_path)
            texto_completo = []
            
            # Extrai texto de parágrafos
            for elem in doc.getElementsByType(text.P):
                texto = teletype.extractText(elem)
                if texto and texto.strip():
                    texto_completo.append(texto.strip())
            
            texto = ' '.join(texto_completo)
            if texto:
                sentencas = nltk.sent_tokenize(texto)
                return [s for s in sentencas if len(s.strip()) > 5]
            return []
        except Exception as e:
            self.estatisticas['erros_processamento'].append(f"ODT {os.path.basename(arquivo_path)}: {e}")
            return []
    
    def processar_audio(self, arquivo_path):
        """Processa arquivos de áudio"""
        if not AUDIO_AVAILABLE:
            print(f"   ⚠️  Processamento de áudio não disponível para {os.path.basename(arquivo_path)}")
            return []
        
        try:
            recognizer = sr.Recognizer()
            
            # Verifica se é MP3 e converte se necessário
            if arquivo_path.lower().endswith('.mp3'):
                try:
                    from pydub import AudioSegment
                    audio = AudioSegment.from_mp3(arquivo_path)
                    wav_path = arquivo_path + '_temp.wav'
                    audio.export(wav_path, format="wav")
                    arquivo_audio = wav_path
                except Exception as e:
                    print(f"   ⚠️  Não foi possível converter MP3: {e}")
                    return []
            else:
                arquivo_audio = arquivo_path
            
            # Transcreve áudio
            with sr.AudioFile(arquivo_audio) as source:
                audio_data = recognizer.record(source)
                texto = recognizer.recognize_google(audio_data, language='pt-BR')
                if texto:
                    sentencas = nltk.sent_tokenize(texto)
                    return [s for s in sentencas if len(s.strip()) > 5]
                return []
                
        except Exception as e:
            print(f"   ⚠️  Erro ao processar áudio: {e}")
            self.estatisticas['erros_processamento'].append(f"ÁUDIO {os.path.basename(arquivo_path)}: {e}")
            return []
        finally:
            # Limpa arquivo temporário se existir
            if 'wav_path' in locals() and os.path.exists(wav_path):
                os.remove(wav_path)
    
    def processar_video(self, arquivo_path):
        """Processa arquivos de vídeo"""
        if not VIDEO_AVAILABLE:
            print(f"   ⚠️  Processamento de vídeo não disponível para {os.path.basename(arquivo_path)}")
            return []
        
        try:
            # Extrai áudio do vídeo
            video = VideoFileClip(arquivo_path)
            audio_path = arquivo_path + "_audio.wav"
            video.audio.write_audiofile(audio_path, verbose=False, logger=None)
            
            # Processa o áudio extraído
            sentencas = self.processar_audio(audio_path)
            
            return sentencas
            
        except Exception as e:
            print(f"   ⚠️  Erro ao processar vídeo: {e}")
            self.estatisticas['erros_processamento'].append(f"VÍDEO {os.path.basename(arquivo_path)}: {e}")
            return []
        finally:
            # Limpa arquivo temporário se existir
            if 'audio_path' in locals() and os.path.exists(audio_path):
                os.remove(audio_path)
    
    def processar_arquivo(self, arquivo_path):
        """Processa qualquer tipo de arquivo"""
        extensao = os.path.splitext(arquivo_path)[1].lower()
        nome_arquivo = os.path.basename(arquivo_path)
        
        mapeamento_processadores = {
            '.txt': self.processar_txt,
            '.docx': self.processar_docx,
            '.doc': self.processar_docx,
            '.pdf': self.processar_pdf,
            '.xlsx': self.processar_xlsx,
            '.xls': self.processar_xlsx,
            '.pptx': self.processar_pptx,
            '.ppt': self.processar_pptx,
            '.odt': self.processar_odt,
            '.rtf': self.processar_txt,  # RTF como texto simples por enquanto
            '.mp3': self.processar_audio,
            '.wav': self.processar_audio,
            '.aac': self.processar_audio,
            '.wma': self.processar_audio,
            '.ogg': self.processar_audio,
            '.mp4': self.processar_video,
            '.avi': self.processar_video,
            '.mov': self.processar_video,
            '.mkv': self.processar_video,
            '.wmv': self.processar_video,
        }
        
        processador = mapeamento_processadores.get(extensao)
        if processador:
            print(f"   🔄 Processando {extensao.upper()}: {nome_arquivo}")
            return processador(arquivo_path)
        else:
            print(f"   ❌ Formato não suportado: {extensao}")
            return []
    
    def aprender_tudo(self, pasta_dados="dados"):
        """Aprende com TODOS os arquivos"""
        print("🧠 Adam Multimodal: Iniciando aprendizado completo...")
        
        # Reset
        self.respostas = []
        self.fontes_conhecimento = {}
        self.estatisticas.update({
            'total_sentencas': 0,
            'total_arquivos': 0,
            'arquivos_processados': [],
            'tipos_processados': {},
            'erros_processamento': []
        })
        
        if not os.path.exists(pasta_dados):
            os.makedirs(pasta_dados)
            print(f"📁 Pasta '{pasta_dados}' criada. Adicione arquivos!")
            return False
        
        # Encontra TODOS os arquivos
        arquivos_encontrados = glob.glob(os.path.join(pasta_dados, "*.*"))
        arquivos_encontrados = [f for f in arquivos_encontrados if os.path.isfile(f)]
        
        if not arquivos_encontrados:
            print("❌ Nenhum arquivo encontrado!")
            return False
        
        print(f"📚 Encontrados {len(arquivos_encontrados)} arquivos:")
        
        total_sentencas = 0
        arquivos_sucesso = 0
        
        for arquivo_path in arquivos_encontrados:
            nome_arquivo = os.path.basename(arquivo_path)
            extensao = os.path.splitext(arquivo_path)[1].lower()
            
            sentencas = self.processar_arquivo(arquivo_path)
            
            if sentencas:
                for sentenca in sentencas:
                    if len(sentenca.strip()) > 5:
                        sentenca_processada = self.preprocessar_texto(sentenca)
                        self.respostas.append(sentenca_processada)
                        self.fontes_conhecimento[len(self.respostas) - 1] = nome_arquivo
                
                total_sentencas += len(sentencas)
                arquivos_sucesso += 1
                self.estatisticas['arquivos_processados'].append(nome_arquivo)
                
                # Estatísticas por tipo
                if extensao in self.estatisticas['tipos_processados']:
                    self.estatisticas['tipos_processados'][extensao] += 1
                else:
                    self.estatisticas['tipos_processados'][extensao] = 1
                
                print(f"   ✅ {nome_arquivo}: {len(sentencas)} sentenças")
            else:
                print(f"   ❌ {nome_arquivo}: 0 sentenças (não processado)")
        
        self.estatisticas['total_sentencas'] = total_sentencas
        self.estatisticas['total_arquivos'] = arquivos_sucesso
        self.estatisticas['ultima_atualizacao'] = np.datetime64('now')
        
        print(f"\n🎉 Adam aprendeu com {arquivos_sucesso}/{len(arquivos_encontrados)} arquivos!")
        print(f"📊 Total: {total_sentencas} sentenças")
        
        # Estatísticas
        if self.estatisticas['tipos_processados']:
            print("📈 Tipos processados:")
            for tipo, qtd in self.estatisticas['tipos_processados'].items():
                print(f"   {tipo}: {qtd} arquivos")
        
        if self.estatisticas['erros_processamento']:
            print(f"⚠️  {len(self.estatisticas['erros_processamento'])} erros de processamento")
        
        return arquivos_sucesso > 0
    
    def treinar_modelo(self):
        """Treina o modelo"""
        if not self.respostas:
            print("❌ Nenhum conhecimento para treinar!")
            return False
        
        print("⚡ Treinando modelo...")
        try:
            self.caracteristicas = self.vectorizer.fit_transform(self.respostas)
            print("✅ Modelo treinado!")
            return True
        except Exception as e:
            print(f"❌ Erro no treinamento: {e}")
            return False
    
    def responder_pergunta(self, pergunta):
        """Responde perguntas"""
        if self.caracteristicas is None:
            return "🤖 Adam: Ainda não fui treinada!"
        
        pergunta_processada = self.preprocessar_texto(pergunta)
        if not pergunta_processada:
            return "🤖 Adam: Por favor, faça uma pergunta mais clara."
        
        try:
            pergunta_vetor = self.vectorizer.transform([pergunta_processada])
            similaridades = cosine_similarity(pergunta_vetor, self.caracteristicas).flatten()
            
            if len(similaridades) == 0:
                return "🤖 Adam: Ainda não tenho conhecimento suficiente."
            
            indice_melhor = similaridades.argmax()
            melhor_similaridade = similaridades[indice_melhor]
            
            if melhor_similaridade < 0.1:
                respostas = [
                    "🤔 Adam: Não encontrei informações sobre isso.",
                    "📚 Adam: Isso não está no meu conhecimento atual.",
                    "💡 Adam: Ainda não aprendi sobre isso.",
                    "🎯 Adam: Hmm, não tenho dados sobre isso."
                ]
                return random.choice(respostas)
            
            fonte = self.fontes_conhecimento.get(indice_melhor, "Conhecimento Geral")
            resposta = self.respostas[indice_melhor]
            
            return f"🧠 Adam (aprendi de {fonte}): {resposta.capitalize()}"
            
        except Exception as e:
            return f"🤖 Adam: Erro ao processar pergunta: {str(e)}"
    
    def adicionar_conhecimento(self, novo_texto, fonte="Manual"):
        """Adiciona conhecimento manual"""
        sentencas = nltk.sent_tokenize(str(novo_texto))
        adicionadas = 0
        
        for sentenca in sentencas:
            if len(sentenca.strip()) > 5:
                sentenca_processada = self.preprocessar_texto(sentenca)
                self.respostas.append(sentenca_processada)
                self.fontes_conhecimento[len(self.respostas) - 1] = fonte
                adicionadas += 1
        
        if adicionadas > 0:
            self.estatisticas['total_sentencas'] = len(self.respostas)
            self.estatisticas['ultima_atualizacao'] = np.datetime64('now')
            self.treinar_modelo()
            print(f"✅ Adam: Aprendi {adicionadas} novas sentenças!")
    
    def get_estatisticas(self):
        """Retorna estatísticas"""
        return {
            'nome': 'Adam - IA Multimodal',
            'status': 'Operacional',
            'total_sentencas': self.estatisticas['total_sentencas'],
            'total_arquivos': self.estatisticas['total_arquivos'],
            'arquivos_processados': self.estatisticas['arquivos_processados'],
            'tipos_processados': self.estatisticas['tipos_processados'],
            'modelo_treinado': self.caracteristicas is not None,
            'versao': '2.0 - Multimodal',
            'erros': len(self.estatisticas['erros_processamento'])
        }